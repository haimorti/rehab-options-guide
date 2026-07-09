"""
Script to generate a professional Hebrew RTL PowerPoint presentation
about tuition reimbursement (שכר לימוד) for students.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.enum.dml import MSO_THEME_COLOR
import copy

# ── Palette ──────────────────────────────────────────────────────────────────
PRIMARY      = RGBColor(0x1A, 0x56, 0xDB)   # deep blue
PRIMARY_DARK = RGBColor(0x10, 0x3A, 0x9E)   # darker blue (header bar)
PRIMARY_LIGHT= RGBColor(0xEB, 0xF0, 0xFF)   # very light blue bg
ACCENT       = RGBColor(0xF5, 0xA6, 0x23)   # amber accent
SUCCESS      = RGBColor(0x05, 0x96, 0x69)   # green checkmark
WARNING      = RGBColor(0xDC, 0x26, 0x26)   # red warning
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
DARK         = RGBColor(0x11, 0x18, 0x27)   # near-black text
GRAY         = RGBColor(0x64, 0x74, 0x8B)   # muted text
LIGHT_GRAY   = RGBColor(0xF1, 0xF5, 0xF9)  # card bg

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ── Helpers ───────────────────────────────────────────────────────────────────

def add_rect(slide, left, top, width, height, fill_color, line_color=None, line_width=None):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, color=DARK,
                align=PP_ALIGN.RIGHT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txBox


def add_rtl_textbox(slide, text, left, top, width, height,
                    font_size=18, bold=False, color=DARK,
                    align=PP_ALIGN.RIGHT, wrap=True, italic=False):
    """Textbox with RTL paragraph direction."""
    from pptx.oxml.ns import qn
    from lxml import etree
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    # Set RTL
    pPr = p._p.get_or_add_pPr()
    pPr.set('rtl', '1')
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txBox


def add_rtl_para(tf, text, font_size=16, bold=False, color=DARK,
                 align=PP_ALIGN.RIGHT, space_before=0, italic=False):
    """Add a paragraph to an existing text frame with RTL."""
    from pptx.oxml.ns import qn
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    pPr = p._p.get_or_add_pPr()
    pPr.set('rtl', '1')
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return p


def set_first_para(tf, text, font_size=16, bold=False, color=DARK,
                   align=PP_ALIGN.RIGHT, italic=False):
    """Configure the first (auto-created) paragraph of a text frame."""
    p = tf.paragraphs[0]
    p.alignment = align
    pPr = p._p.get_or_add_pPr()
    pPr.set('rtl', '1')
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"


def add_bullet_item(slide, icon, text, left, top, width,
                    font_size=15, text_color=DARK, icon_color=SUCCESS,
                    bold_text=False):
    """Single line: icon + text side by side (RTL: text right, icon left)."""
    icon_w = Inches(0.35)
    gap = Inches(0.1)
    text_w = width - icon_w - gap

    # text first (right side)
    txBox = slide.shapes.add_textbox(left, top, text_w, Inches(0.38))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    pPr = p._p.get_or_add_pPr()
    pPr.set('rtl', '1')
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold_text
    run.font.color.rgb = text_color
    run.font.name = "Calibri"

    # icon circle (left of text in RTL layout = right margin side)
    ico = slide.shapes.add_textbox(left + text_w + gap, top, icon_w, Inches(0.38))
    tf2 = ico.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = icon
    run2.font.size = Pt(font_size + 2)
    run2.font.color.rgb = icon_color
    run2.font.name = "Segoe UI Emoji"

    return txBox


def add_slide_header(slide, title, subtitle=None,
                     bg_color=PRIMARY_DARK, text_color=WHITE):
    """Top header bar with title and optional subtitle."""
    bar_h = Inches(1.35) if subtitle else Inches(1.0)
    add_rect(slide, 0, 0, SLIDE_W, bar_h, bg_color)

    # Decorative right accent stripe
    add_rect(slide, SLIDE_W - Inches(0.25), 0, Inches(0.25), bar_h, ACCENT)

    # Title
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), SLIDE_W - Inches(1.0), Inches(0.65))
    tf = tx.text_frame
    tf.word_wrap = False
    set_first_para(tf, title, font_size=28, bold=True, color=text_color)

    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.5), Inches(0.75), SLIDE_W - Inches(1.0), Inches(0.45))
        tf2 = sub.text_frame
        tf2.word_wrap = False
        set_first_para(tf2, subtitle, font_size=15, bold=False,
                       color=RGBColor(0xBF, 0xDB, 0xFF), italic=False)

    return bar_h


def number_badge(slide, number, cx, cy, radius=Inches(0.22), bg=PRIMARY, fg=WHITE):
    """Draw a filled circle with a number (uses a rounded rectangle as approximation)."""
    from pptx.util import Emu
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    # Use an oval shape
    shape = slide.shapes.add_shape(
        9,  # oval
        cx - radius, cy - radius, radius * 2, radius * 2
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg
    shape.line.fill.background()

    tx = slide.shapes.add_textbox(cx - radius, cy - radius, radius * 2, radius * 2)
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = str(number)
    run.font.size = Pt(12)
    run.font.bold = True
    run.font.color.rgb = fg
    run.font.name = "Calibri"


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE BUILDERS
# ══════════════════════════════════════════════════════════════════════════════

def build_title_slide(prs):
    """Slide 1 – Title / cover."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Full bg
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, PRIMARY_DARK)

    # Bottom accent strip
    add_rect(slide, 0, SLIDE_H - Inches(0.3), SLIDE_W, Inches(0.3), ACCENT)

    # Decorative large circle (top-right)
    r = Inches(3.5)
    circ = slide.shapes.add_shape(9, SLIDE_W - r * 0.7, -r * 0.4, r, r)
    circ.fill.solid()
    circ.fill.fore_color.rgb = RGBColor(0x1E, 0x65, 0xF0)
    circ.line.fill.background()

    # Decorative small circle (bottom-left)
    r2 = Inches(2.0)
    circ2 = slide.shapes.add_shape(9, -r2 * 0.5, SLIDE_H - r2 * 0.7, r2, r2)
    circ2.fill.solid()
    circ2.fill.fore_color.rgb = RGBColor(0x1E, 0x65, 0xF0)
    circ2.line.fill.background()

    # Main title
    tx_main = slide.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.3), Inches(1.5))
    tf = tx_main.text_frame
    set_first_para(tf, "שכר לימוד", font_size=52, bold=True,
                   color=WHITE, align=PP_ALIGN.CENTER)

    # Subtitle
    tx_sub = slide.shapes.add_textbox(Inches(1.5), Inches(3.4), Inches(10.3), Inches(0.8))
    tf2 = tx_sub.text_frame
    set_first_para(tf2, "כל מה שצריך לדעת על ההחזר שמגיע לך",
                   font_size=22, bold=False,
                   color=RGBColor(0xBF, 0xDB, 0xFF), align=PP_ALIGN.CENTER)

    # Divider line
    ln = slide.shapes.add_shape(1, Inches(4.5), Inches(4.3), Inches(4.3), Pt(2))
    ln.fill.solid()
    ln.fill.fore_color.rgb = ACCENT
    ln.line.fill.background()

    # Source label
    tx3 = slide.shapes.add_textbox(Inches(1.0), Inches(4.7), Inches(11.3), Inches(0.5))
    tf3 = tx3.text_frame
    set_first_para(tf3, "מיועד לסטודנטים המקבלים תמיכה בשיקום מקצועי",
                   font_size=14, color=RGBColor(0x94, 0xA3, 0xB8),
                   align=PP_ALIGN.CENTER)

    return slide


def build_overview_slide(prs):
    """Slide 2 – Overview / summary."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, LIGHT_GRAY)
    bar_h = add_slide_header(slide, "בקצרה – מה שחשוב לדעת")

    # 4 info cards
    cards = [
        ("💰", "תקרת ההחזר", "עד 13,079 ₪\nלשנת לימודים (2025)"),
        ("🏦", "אופן התשלום", "זיכוי ישיר לחשבון הבנק\nכנגד קבלות בלבד"),
        ("📅", "שתי פעימות", "פעימה א' – סמסטר א'\nפעימה ב' – אחרי אישור סמסטר ב'"),
        ("⚖️", "הפרש מעל התקרה", "אם שכר הלימוד גבוה מהתקרה –\nתישאו בהפרש באופן עצמאי"),
    ]

    card_w = Inches(2.95)
    card_h = Inches(2.5)
    gap = Inches(0.25)
    start_x = Inches(0.35)
    top_y = bar_h + Inches(0.35)

    for i, (icon, title, body) in enumerate(cards):
        cx = start_x + i * (card_w + gap)

        # Card bg
        card = add_rect(slide, cx, top_y, card_w, card_h, WHITE,
                        line_color=RGBColor(0xCB, 0xD5, 0xE1), line_width=Pt(1))
        card.shadow.inherit = False

        # Top accent
        add_rect(slide, cx, top_y, card_w, Inches(0.07), PRIMARY)

        # Icon
        ico_box = slide.shapes.add_textbox(cx, top_y + Inches(0.15), card_w, Inches(0.6))
        tf_ico = ico_box.text_frame
        p = tf_ico.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = icon
        run.font.size = Pt(28)
        run.font.name = "Segoe UI Emoji"

        # Title
        ttl = slide.shapes.add_textbox(cx + Inches(0.15), top_y + Inches(0.8), card_w - Inches(0.3), Inches(0.4))
        tf_ttl = ttl.text_frame
        set_first_para(tf_ttl, title, font_size=14, bold=True, color=PRIMARY_DARK, align=PP_ALIGN.CENTER)

        # Body
        bdy = slide.shapes.add_textbox(cx + Inches(0.15), top_y + Inches(1.25), card_w - Inches(0.3), Inches(1.1))
        tf_bdy = bdy.text_frame
        tf_bdy.word_wrap = True
        set_first_para(tf_bdy, body, font_size=13, color=GRAY, align=PP_ALIGN.CENTER)

    # Bottom note
    note = slide.shapes.add_textbox(Inches(0.5), top_y + card_h + Inches(0.3),
                                    SLIDE_W - Inches(1.0), Inches(0.55))
    tf_note = note.text_frame
    set_first_para(tf_note,
                   "ההחזר אינו מועבר ישירות למוסד הלימודים – יש לשלם בעצמך ולאחר מכן להגיש קבלה.",
                   font_size=13, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)

    return slide


def build_ceiling_slide(prs):
    """Slide 3 – תקרת ההחזר."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, LIGHT_GRAY)
    bar_h = add_slide_header(slide, "תקרת ההחזר", "כמה כסף מגיע לך ואיך מחשבים את זה")

    content_top = bar_h + Inches(0.25)

    # ── Left column: rules ──────────────────────────────────────────────
    col_l = Inches(0.4)
    col_w = Inches(6.5)

    rules = [
        ("✅", "תקרה שנתית: 13,079 ₪ (נכון לשנת 2025)", SUCCESS),
        ("✅", "ההחזר מועבר ישירות לחשבון הבנק שלך", SUCCESS),
        ("✅", "נדרשת קבלה על תשלום שבוצע בפועל", SUCCESS),
        ("⚠️", "הלימוד לא מגיע ישירות למוסד – אתה משלם קודם", WARNING),
        ("⚠️", "שכר לימוד מעל התקרה? ההפרש על חשבונך", WARNING),
    ]

    row_h = Inches(0.55)
    for i, (icon, text, col) in enumerate(rules):
        y = content_top + Inches(0.1) + i * row_h
        bg = RGBColor(0xF0, 0xFD, 0xF4) if col == SUCCESS else RGBColor(0xFF, 0xF5, 0xF5)
        add_rect(slide, col_l, y, col_w, row_h - Inches(0.06), bg,
                 line_color=RGBColor(0xD1, 0xD5, 0xDB), line_width=Pt(0.5))
        add_bullet_item(slide, icon, text,
                        col_l + Inches(0.12), y + Inches(0.08),
                        col_w - Inches(0.25),
                        font_size=14, text_color=DARK, icon_color=col)

    # ── Right column: example box ────────────────────────────────────────
    ex_l = Inches(7.2)
    ex_w = Inches(5.7)
    ex_h = Inches(3.0)

    add_rect(slide, ex_l, content_top, ex_w, ex_h, WHITE,
             line_color=PRIMARY, line_width=Pt(2))
    add_rect(slide, ex_l, content_top, ex_w, Inches(0.45), PRIMARY)

    # Header of example
    hdr = slide.shapes.add_textbox(ex_l + Inches(0.15), content_top + Inches(0.05),
                                   ex_w - Inches(0.3), Inches(0.38))
    set_first_para(hdr.text_frame, "דוגמה להחזר יחסי",
                   font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    ex_body = slide.shapes.add_textbox(ex_l + Inches(0.25), content_top + Inches(0.6),
                                       ex_w - Inches(0.5), ex_h - Inches(0.75))
    tf_ex = ex_body.text_frame
    tf_ex.word_wrap = True
    set_first_para(tf_ex,
                   "נניח ששכר הלימוד השנתי שלך הוא 26,158 ₪ –",
                   font_size=13, color=DARK, align=PP_ALIGN.RIGHT)
    add_rtl_para(tf_ex, "כפול מהתקרה (13,079 ₪).",
                 font_size=13, color=DARK, space_before=2)
    add_rtl_para(tf_ex, "", font_size=6, space_before=0)
    add_rtl_para(tf_ex,
                 "במצב כזה תקבל החזר של 50% מכל קבלה,",
                 font_size=13, color=PRIMARY_DARK, bold=True, space_before=2)
    add_rtl_para(tf_ex,
                 "עד לניצול מלא של 13,079 ₪.",
                 font_size=13, color=PRIMARY_DARK, bold=True, space_before=2)
    add_rtl_para(tf_ex, "", font_size=6)
    add_rtl_para(tf_ex,
                 "כלומר – ההחזר מחושב יחסית לפי התקרה.",
                 font_size=12, color=GRAY, italic=True, space_before=2)

    # Divider formula
    form_top = content_top + ex_h + Inches(0.2)
    add_rect(slide, ex_l, form_top, ex_w, Inches(1.45),
             PRIMARY_LIGHT, line_color=PRIMARY, line_width=Pt(1))
    formula = slide.shapes.add_textbox(ex_l + Inches(0.2), form_top + Inches(0.12),
                                       ex_w - Inches(0.4), Inches(1.2))
    tf_f = formula.text_frame
    tf_f.word_wrap = True
    set_first_para(tf_f, "📐 נוסחת החזר יחסי", font_size=13, bold=True,
                   color=PRIMARY_DARK, align=PP_ALIGN.RIGHT)
    add_rtl_para(tf_f,
                 "אחוז ההחזר = (תקרה ÷ שכר לימוד בפועל) × 100",
                 font_size=12, color=PRIMARY_DARK, space_before=4)
    add_rtl_para(tf_f,
                 "הסכום שתקבל = אחוז ההחזר × סכום כל קבלה",
                 font_size=12, color=PRIMARY_DARK, space_before=2)

    return slide


def build_payment_timing_slide(prs):
    """Slide 4 – מתי מקבלים את ההחזר (פעימות)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, LIGHT_GRAY)
    bar_h = add_slide_header(slide, "מתי מקבלים את ההחזר?",
                             "אישור הזכאות לכל סמסטר נבדק בנפרד – ולכן ההחזר משולם בשתי פעימות")

    top = bar_h + Inches(0.3)

    # ── Two step cards (horizontal) ─────────────────────────────────────
    card_w = Inches(5.6)
    card_h = Inches(2.6)
    gap_between = Inches(0.6)
    left1 = Inches(0.6)
    left2 = left1 + card_w + gap_between

    def payment_card(slide, left, top, num, title, body, badge_color):
        add_rect(slide, left, top, card_w, card_h, WHITE,
                 line_color=badge_color, line_width=Pt(2))
        # badge strip
        add_rect(slide, left, top, card_w, Inches(0.55), badge_color)
        # number
        num_box = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.06),
                                           Inches(0.45), Inches(0.45))
        set_first_para(num_box.text_frame, str(num), font_size=18, bold=True,
                       color=WHITE, align=PP_ALIGN.CENTER)
        # card title
        t = slide.shapes.add_textbox(left + Inches(0.6), top + Inches(0.07),
                                     card_w - Inches(0.75), Inches(0.42))
        set_first_para(t.text_frame, title, font_size=16, bold=True,
                       color=WHITE, align=PP_ALIGN.RIGHT)
        # body
        b = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.7),
                                     card_w - Inches(0.4), card_h - Inches(0.85))
        tf_b = b.text_frame
        tf_b.word_wrap = True
        set_first_para(tf_b, body, font_size=14, color=DARK, align=PP_ALIGN.RIGHT)

    payment_card(slide, left1, top, 1, "פעימה ראשונה – סמסטר א'",
                 "מחצית מסכום הזכאות השנתית תשולם לחשבונך בתחילת סמסטר א'.\n\nהתשלום מותנה באישור הזכאות לסמסטר זה.",
                 PRIMARY)

    payment_card(slide, left2, top, 2, "פעימה שנייה – סמסטר ב'",
                 "המחצית השנייה תשולם לאחר אישור סמסטר ב'.\n\nהאישור ניתן רק לאחר בחינת התקדמותך בלימודים.",
                 RGBColor(0x05, 0x7A, 0x55))

    # ── Bottom info box ──────────────────────────────────────────────────
    info_top = top + card_h + Inches(0.3)
    add_rect(slide, Inches(0.6), info_top, SLIDE_W - Inches(1.2), Inches(1.4),
             PRIMARY_LIGHT, line_color=PRIMARY, line_width=Pt(1))

    info_tx = slide.shapes.add_textbox(Inches(0.9), info_top + Inches(0.12),
                                       SLIDE_W - Inches(1.8), Inches(1.15))
    tf_info = info_tx.text_frame
    tf_info.word_wrap = True
    set_first_para(tf_info, "💡 מקדמה לפני תחילת הלימודים",
                   font_size=14, bold=True, color=PRIMARY_DARK, align=PP_ALIGN.RIGHT)
    add_rtl_para(tf_info,
                 "במקרים מסוימים, כאשר שולמה מקדמה, ניתן לקבל אותה לפני מועד תחילת הלימודים הרשמי.",
                 font_size=13, color=DARK, space_before=4)

    return slide


def build_scholarships_slide(prs):
    """Slide 5 – השפעת מלגות, הנחות ופטורים."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, LIGHT_GRAY)
    bar_h = add_slide_header(slide, "מלגות, הנחות ופטורים",
                             "לא כל מלגה משפיעה על ההחזר – הנה מה שחשוב לדעת")

    top = bar_h + Inches(0.35)
    col_w = Inches(5.8)
    col_h = Inches(3.2)
    gap = Inches(0.55)
    left1 = Inches(0.5)
    left2 = left1 + col_w + gap

    # ── Card A: מלגות שמנוכות ────────────────────────────────────────────
    add_rect(slide, left1, top, col_w, col_h, WHITE,
             line_color=WARNING, line_width=Pt(2))
    add_rect(slide, left1, top, col_w, Inches(0.5), WARNING)

    hdr_a = slide.shapes.add_textbox(left1 + Inches(0.2), top + Inches(0.06),
                                     col_w - Inches(0.4), Inches(0.4))
    set_first_para(hdr_a.text_frame, "❌  מלגות שמנוכות מההחזר",
                   font_size=15, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)

    body_a = slide.shapes.add_textbox(left1 + Inches(0.2), top + Inches(0.65),
                                      col_w - Inches(0.4), col_h - Inches(0.8))
    tf_a = body_a.text_frame
    tf_a.word_wrap = True
    set_first_para(tf_a,
                   "מלגות ייעודיות שניתנו במיוחד עבור שכר לימוד –",
                   font_size=13, bold=True, color=DARK, align=PP_ALIGN.RIGHT)
    add_rtl_para(tf_a, "ינוכו מסכום ההחזר שתקבל.", font_size=13, color=DARK, space_before=2)
    add_rtl_para(tf_a, "", font_size=6)
    add_rtl_para(tf_a, "דוגמאות:", font_size=13, bold=True, color=WARNING, space_before=4)
    add_rtl_para(tf_a, '• מלגת "ממדים ללימודים"', font_size=13, color=DARK, space_before=2)
    add_rtl_para(tf_a, "• מלגת משרד הביטחון ללימודים בפריפריה", font_size=13, color=DARK, space_before=2)

    # ── Card B: מלגות שאינן מנוכות ──────────────────────────────────────
    add_rect(slide, left2, top, col_w, col_h, WHITE,
             line_color=SUCCESS, line_width=Pt(2))
    add_rect(slide, left2, top, col_w, Inches(0.5), SUCCESS)

    hdr_b = slide.shapes.add_textbox(left2 + Inches(0.2), top + Inches(0.06),
                                     col_w - Inches(0.4), Inches(0.4))
    set_first_para(hdr_b.text_frame, "✅  מלגות שאינן מנוכות",
                   font_size=15, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)

    body_b = slide.shapes.add_textbox(left2 + Inches(0.2), top + Inches(0.65),
                                      col_w - Inches(0.4), col_h - Inches(0.8))
    tf_b = body_b.text_frame
    tf_b.word_wrap = True
    set_first_para(tf_b,
                   "מלגות שאינן מועברות ישירות לחשבון שכר הלימוד",
                   font_size=13, bold=True, color=DARK, align=PP_ALIGN.RIGHT)
    add_rtl_para(tf_b,
                 "ויכולות לשמש למטרות אחרות –",
                 font_size=13, color=DARK, space_before=2)
    add_rtl_para(tf_b, "אינן נלקחות בחשבון בחישוב ההחזר.", font_size=13, color=SUCCESS,
                 bold=True, space_before=4)
    add_rtl_para(tf_b, "", font_size=6)
    add_rtl_para(tf_b, "כלומר: מלגות כאלה לא יפגעו בסכום שתקבל.",
                 font_size=13, color=GRAY, italic=True, space_before=2)

    # ── Bottom tip ────────────────────────────────────────────────────────
    tip_top = top + col_h + Inches(0.3)
    add_rect(slide, Inches(0.5), tip_top, SLIDE_W - Inches(1.0), Inches(1.1),
             RGBColor(0xFF, 0xF8, 0xE7), line_color=ACCENT, line_width=Pt(1.5))
    tip_tx = slide.shapes.add_textbox(Inches(0.8), tip_top + Inches(0.12),
                                      SLIDE_W - Inches(1.6), Inches(0.85))
    tf_tip = tip_tx.text_frame
    tf_tip.word_wrap = True
    set_first_para(tf_tip,
                   "💡 טיפ: אם קיבלת מלגה – בדוק עם עו\"ס השיקום שלך האם היא משפיעה על גובה ההחזר.",
                   font_size=13, bold=True, color=RGBColor(0x92, 0x40, 0x0E), align=PP_ALIGN.RIGHT)

    return slide


def build_receipt_slide(prs):
    """Slide 6 – הנחיות להגשת קבלה."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, LIGHT_GRAY)
    bar_h = add_slide_header(slide, "כיצד מגישים קבלה?",
                             "לאחר התשלום למוסד הלימודים – שלחו את הקבלה לעו\"ס השיקום דרך אתר הביטוח הלאומי")

    top = bar_h + Inches(0.3)

    # Step flow: 3 steps top
    steps = [
        ("1", "שלמו את שכר הלימוד", "למוסד האקדמי ישירות"),
        ("2", "קבלו קבלה / חשבונית", "מהמוסד האקדמי"),
        ("3", "שלחו לעו\"ס השיקום", "דרך אתר הביטוח הלאומי"),
    ]
    step_w = Inches(3.8)
    step_h = Inches(1.15)
    step_gap = Inches(0.3)
    step_start = Inches(0.5)

    for i, (num, t1, t2) in enumerate(steps):
        sx = step_start + i * (step_w + step_gap)
        add_rect(slide, sx, top, step_w, step_h, WHITE,
                 line_color=PRIMARY, line_width=Pt(1.5))
        add_rect(slide, sx, top, Inches(0.55), step_h, PRIMARY)
        # number
        nb = slide.shapes.add_textbox(sx + Inches(0.05), top + Inches(0.28),
                                      Inches(0.45), Inches(0.5))
        set_first_para(nb.text_frame, num, font_size=18, bold=True,
                       color=WHITE, align=PP_ALIGN.CENTER)
        # text
        tb = slide.shapes.add_textbox(sx + Inches(0.65), top + Inches(0.12),
                                      step_w - Inches(0.75), step_h - Inches(0.2))
        tf_s = tb.text_frame
        tf_s.word_wrap = True
        set_first_para(tf_s, t1, font_size=14, bold=True, color=DARK, align=PP_ALIGN.RIGHT)
        add_rtl_para(tf_s, t2, font_size=12, color=GRAY, space_before=2)

        # Arrow between steps
        if i < len(steps) - 1:
            arr_x = sx + step_w + Inches(0.03)
            arr = slide.shapes.add_textbox(arr_x, top + Inches(0.35), Inches(0.25), Inches(0.45))
            set_first_para(arr.text_frame, "◄", font_size=14, color=PRIMARY, align=PP_ALIGN.CENTER)

    # Checklist
    cl_top = top + step_h + Inches(0.35)
    cl_l = Inches(0.5)
    cl_w = SLIDE_W - Inches(1.0)
    cl_header_h = Inches(0.48)

    add_rect(slide, cl_l, cl_top, cl_w, cl_header_h, PRIMARY)
    cl_hdr = slide.shapes.add_textbox(cl_l + Inches(0.2), cl_top + Inches(0.07),
                                      cl_w - Inches(0.4), cl_header_h - Inches(0.1))
    set_first_para(cl_hdr.text_frame, "🧾  מה חייב להופיע על הקבלה",
                   font_size=15, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)

    items = [
        "שם מלא",
        "מספר תעודת זהות",
        "שם המוסד האקדמי",
        "תאריך הנפקת הקבלה",
        "מספר קבלה / חשבונית",
        "מהות התשלום – שכר לימוד",
    ]

    item_h = Inches(0.42)
    per_row = 3
    rows = [items[i:i+per_row] for i in range(0, len(items), per_row)]
    col_w_item = cl_w / per_row

    for r, row in enumerate(rows):
        for c, item in enumerate(row):
            iy = cl_top + cl_header_h + r * item_h
            ix = cl_l + c * col_w_item
            bg = WHITE if (r + c) % 2 == 0 else RGBColor(0xF8, 0xFA, 0xFF)
            add_rect(slide, ix, iy, col_w_item, item_h, bg,
                     line_color=RGBColor(0xD1, 0xD5, 0xDB), line_width=Pt(0.5))
            add_bullet_item(slide, "✓", item,
                            ix + Inches(0.1), iy + Inches(0.07),
                            col_w_item - Inches(0.2),
                            font_size=13, text_color=DARK, icon_color=SUCCESS)

    # Warning bar
    warn_top = cl_top + cl_header_h + len(rows) * item_h
    add_rect(slide, cl_l, warn_top, cl_w, Inches(0.52),
             RGBColor(0xFF, 0xF0, 0xF0), line_color=WARNING, line_width=Pt(1))
    warn_tx = slide.shapes.add_textbox(cl_l + Inches(0.2), warn_top + Inches(0.09),
                                       cl_w - Inches(0.4), Inches(0.38))
    set_first_para(warn_tx.text_frame,
                   "⚠️  קבלה החסרה אחד או יותר מהפרטים הנדרשים תידחה – אנא וודאו שהכל מלא לפני ההגשה.",
                   font_size=13, bold=True, color=WARNING, align=PP_ALIGN.RIGHT)

    return slide


def build_processing_time_slide(prs):
    """Slide 7 – תוך כמה זמן מתקבל ההחזר."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, LIGHT_GRAY)
    bar_h = add_slide_header(slide, "תוך כמה זמן מגיע הכסף?")

    top = bar_h + Inches(0.35)

    # Big central time display
    center_box = Inches(4.5)
    center_l = (SLIDE_W - center_box) / 2
    add_rect(slide, center_l, top, center_box, Inches(2.4), WHITE,
             line_color=PRIMARY, line_width=Pt(2))
    add_rect(slide, center_l, top, center_box, Inches(0.12), PRIMARY)

    ico = slide.shapes.add_textbox(center_l, top + Inches(0.2), center_box, Inches(0.7))
    set_first_para(ico.text_frame, "⏱️", font_size=36, align=PP_ALIGN.CENTER)

    big_num = slide.shapes.add_textbox(center_l, top + Inches(0.95), center_box, Inches(0.8))
    set_first_para(big_num.text_frame, "~10 ימי עבודה", font_size=28, bold=True,
                   color=PRIMARY_DARK, align=PP_ALIGN.CENTER)

    sub_lbl = slide.shapes.add_textbox(center_l, top + Inches(1.75), center_box, Inches(0.5))
    set_first_para(sub_lbl.text_frame, "מרגע הגשת הקבלה ועד קבלת הזיכוי לחשבון",
                   font_size=13, color=GRAY, align=PP_ALIGN.CENTER)

    # Notes below
    note_top = top + Inches(2.7)

    notes = [
        ("✅", "ברוב המקרים ההחזר מתקבל תוך כ-10 ימי עבודה.", SUCCESS),
        ("⚠️", "ייתכנו עיכובים בזמני עומס חריג.", WARNING),
        ("⚠️", "ייתכנו עיכובים בתקופת חגים וחופשות.", WARNING),
        ("💡", "מומלץ לשמור עותק של הקבלה לאחר ההגשה.", ACCENT),
    ]

    note_h = Inches(0.55)
    note_w = Inches(10.0)
    note_l = (SLIDE_W - note_w) / 2

    for i, (icon, text, col) in enumerate(notes):
        ny = note_top + i * (note_h + Inches(0.06))
        bg = RGBColor(0xF0, 0xFD, 0xF4) if col == SUCCESS else \
             RGBColor(0xFF, 0xF5, 0xF5) if col == WARNING else \
             RGBColor(0xFF, 0xFB, 0xEB)
        add_rect(slide, note_l, ny, note_w, note_h, bg,
                 line_color=col, line_width=Pt(1))
        add_bullet_item(slide, icon, text, note_l + Inches(0.15), ny + Inches(0.1),
                        note_w - Inches(0.3), font_size=14,
                        text_color=DARK, icon_color=col)

    return slide


def build_conditions_slide(prs):
    """Slide 8 – תנאים להמשך קבלת התמיכה."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, LIGHT_GRAY)
    bar_h = add_slide_header(slide, "תנאים להמשך קבלת התמיכה",
                             "המשך הזכאות מותנה בהתקדמות בלימודים ובהגשת מסמכים בזמן")

    top = bar_h + Inches(0.35)

    # Warning banner
    add_rect(slide, Inches(0.5), top, SLIDE_W - Inches(1.0), Inches(1.05),
             RGBColor(0xFF, 0xF0, 0xEB), line_color=WARNING, line_width=Pt(2))
    warn_ico = slide.shapes.add_textbox(Inches(0.7), top + Inches(0.22),
                                        Inches(0.7), Inches(0.6))
    set_first_para(warn_ico.text_frame, "⚠️", font_size=28, align=PP_ALIGN.CENTER)
    warn_tx = slide.shapes.add_textbox(Inches(1.4), top + Inches(0.12),
                                       SLIDE_W - Inches(2.2), Inches(0.8))
    tf_w = warn_tx.text_frame
    tf_w.word_wrap = True
    set_first_para(tf_w,
                   "בסיום כל סמסטר עליך להגיש מסמכים – אחרת התמיכה עלולה להיפסק.",
                   font_size=16, bold=True, color=WARNING, align=PP_ALIGN.RIGHT)

    # Two required doc cards
    docs_top = top + Inches(1.25)
    doc_w = Inches(5.7)
    doc_h = Inches(2.5)
    doc_gap = Inches(0.6)
    doc_l1 = Inches(0.6)
    doc_l2 = doc_l1 + doc_w + doc_gap

    def doc_card(slide, left, top, icon, title, desc, border_color):
        add_rect(slide, left, top, doc_w, doc_h, WHITE,
                 line_color=border_color, line_width=Pt(2))
        ico_b = slide.shapes.add_textbox(left, top + Inches(0.2), doc_w, Inches(0.65))
        set_first_para(ico_b.text_frame, icon, font_size=32, align=PP_ALIGN.CENTER)
        t = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.9),
                                     doc_w - Inches(0.4), Inches(0.45))
        set_first_para(t.text_frame, title, font_size=16, bold=True,
                       color=border_color, align=PP_ALIGN.CENTER)
        d = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(1.4),
                                     doc_w - Inches(0.4), doc_h - Inches(1.55))
        tf_d = d.text_frame
        tf_d.word_wrap = True
        set_first_para(tf_d, desc, font_size=13, color=DARK, align=PP_ALIGN.CENTER)

    doc_card(slide, doc_l1, docs_top, "📋", "גיליון ציונים",
             "ציוני הסמסטר המסתיים –\nהוכחה להתקדמות בלימודים.",
             PRIMARY)
    doc_card(slide, doc_l2, docs_top, "📅", "מערכת שעות מעודכנת",
             "לסמסטר הבא –\nהוכחה לרישום לקורסים.",
             RGBColor(0x05, 0x7A, 0x55))

    # Bottom condition note
    cond_top = docs_top + doc_h + Inches(0.3)
    add_rect(slide, Inches(0.5), cond_top, SLIDE_W - Inches(1.0), Inches(0.95),
             PRIMARY_LIGHT, line_color=PRIMARY, line_width=Pt(1))
    cond_tx = slide.shapes.add_textbox(Inches(0.8), cond_top + Inches(0.12),
                                       SLIDE_W - Inches(1.6), Inches(0.7))
    tf_c = cond_tx.text_frame
    tf_c.word_wrap = True
    set_first_para(tf_c,
                   "המשך קבלת כל התמיכות (שכר לימוד, הסעות, ציוד וכו') מותנה "
                   "בעמידה בדרישות הלימודים ובהצלחה במבחני הסמסטר הקודם.",
                   font_size=13, bold=True, color=PRIMARY_DARK, align=PP_ALIGN.RIGHT)

    return slide


def build_summary_slide(prs):
    """Slide 9 – Summary / closing."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, PRIMARY_DARK)
    add_rect(slide, 0, SLIDE_H - Inches(0.3), SLIDE_W, Inches(0.3), ACCENT)

    # Decorative circles
    r = Inches(3.0)
    c = slide.shapes.add_shape(9, SLIDE_W - r * 0.6, -r * 0.3, r, r)
    c.fill.solid()
    c.fill.fore_color.rgb = RGBColor(0x1E, 0x65, 0xF0)
    c.line.fill.background()

    title = slide.shapes.add_textbox(Inches(1.0), Inches(0.5), SLIDE_W - Inches(2.0), Inches(0.9))
    set_first_para(title.text_frame, "סיכום – נקודות מפתח",
                   font_size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # 6 key points in 2 columns
    points = [
        ("💰", "תקרה: 13,079 ₪ לשנה (2025)"),
        ("🏦", "זיכוי לחשבון הבנק – לא ישירות למוסד"),
        ("📅", "2 פעימות: סמסטר א' + סמסטר ב'"),
        ("🧾", "חובה להגיש קבלה עם כל הפרטים"),
        ("⏱️", "~10 ימי עבודה לקבלת ההחזר"),
        ("📋", "בסיום כל סמסטר – הגשת ציונים ומערכת שעות"),
    ]

    pt_w = Inches(5.5)
    pt_h = Inches(0.72)
    pt_gap_v = Inches(0.15)
    pt_gap_h = Inches(0.5)
    pt_l1 = Inches(0.8)
    pt_l2 = pt_l1 + pt_w + pt_gap_h
    pt_top = Inches(1.6)

    for i, (icon, text) in enumerate(points):
        col = i % 2
        row = i // 2
        px = pt_l1 if col == 0 else pt_l2
        py = pt_top + row * (pt_h + pt_gap_v)

        add_rect(slide, px, py, pt_w, pt_h,
                 RGBColor(0x1C, 0x3F, 0xAA),
                 line_color=RGBColor(0x2D, 0x5C, 0xE0), line_width=Pt(1))

        ico_b = slide.shapes.add_textbox(px + pt_w - Inches(0.7), py + Inches(0.12),
                                         Inches(0.55), Inches(0.5))
        set_first_para(ico_b.text_frame, icon, font_size=22, align=PP_ALIGN.CENTER)

        tx_b = slide.shapes.add_textbox(px + Inches(0.15), py + Inches(0.15),
                                        pt_w - Inches(0.9), pt_h - Inches(0.25))
        tf_pt = tx_b.text_frame
        tf_pt.word_wrap = True
        set_first_para(tf_pt, text, font_size=14, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)

    # Closing line
    close_top = pt_top + 3 * (pt_h + pt_gap_v) + Inches(0.25)
    cl = slide.shapes.add_textbox(Inches(1.0), close_top, SLIDE_W - Inches(2.0), Inches(0.55))
    set_first_para(cl.text_frame,
                   "לשאלות נוספות – פנו לעו\"ס השיקום שלכם",
                   font_size=15, color=RGBColor(0xBF, 0xDB, 0xFF), align=PP_ALIGN.CENTER)

    return slide


# ══════════════════════════════════════════════════════════════════════════════
#  MAIN
# ══════════════════════════════════════════════════════════════════════════════

def main():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    build_title_slide(prs)
    build_overview_slide(prs)
    build_ceiling_slide(prs)
    build_payment_timing_slide(prs)
    build_scholarships_slide(prs)
    build_receipt_slide(prs)
    build_processing_time_slide(prs)
    build_conditions_slide(prs)
    build_summary_slide(prs)

    out = "/home/user/shikum-accessibility/tuition_presentation.pptx"
    prs.save(out)
    print(f"Saved: {out}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
